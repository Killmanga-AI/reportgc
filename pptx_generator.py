from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime

class PPTXGenerator:
    def __init__(self, master_pptx: Path = None):
        self.prs = Presentation(str(master_pptx)) if master_pptx and master_pptx.exists() else Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _get_color(self, grade: str) -> RGBColor:
        colors = {
            'A': RGBColor(40, 167, 69), 'B': RGBColor(108, 117, 125),
            'C': RGBColor(255, 193, 7), 'D': RGBColor(253, 126, 20),
            'F': RGBColor(220, 53, 69)
        }
        return colors.get(grade, RGBColor(0, 0, 0))

    def _get_risk_color(self, risk_level: str) -> RGBColor:
        """Get color for risk level badges."""
        colors = {
            'FULL_TABLE_SCAN': RGBColor(220, 53, 69),   # Critical - Red
            'INDEX_RANGE_SCAN': RGBColor(253, 126, 20), # High - Orange
            'NESTED_LOOP': RGBColor(255, 193, 7),       # Medium - Yellow
            'SEQUENTIAL_READ': RGBColor(108, 117, 125)  # Low - Gray
        }
        return colors.get(risk_level, RGBColor(108, 117, 125))

    def _ensure_data_structure(self, data: dict) -> dict:
        """Sanitizes input data for slide stability."""
        data.setdefault('grade', 'F')
        data.setdefault('generated_at', datetime.now().strftime('%Y-%m-%d %H:%M'))
        data.setdefault('summary', {'total_findings': 0})
        
        ep = data.setdefault('execution_plan', {})
        
        # Support both 3-tier (legacy) and 4-tier (new) engine outputs
        for section in ['full_table_scans', 'index_scans', 'nested_loops', 'low_priority']:
            ep.setdefault(section, {'count': 0, 'estimated_hours': 0, 'items': []})
        
        # Backward compatibility: if no nested_loops but engine has medium count
        if ep['nested_loops']['count'] == 0 and data.get('summary', {}).get('medium', 0) > 0:
            # Medium findings were lumped into low_priority, split them visually
            pass  # Keep as-is, will show in matrix
        
        return data

    def generate_pptx(self, data: dict, output_path: str):
        data = self._ensure_data_structure(data)
        self._add_title_slide(data)
        self._add_matrix_slide(data)
        self._add_critical_detail_slide(data)
        self._add_high_detail_slide(data)  # New: High severity details
        self._add_roadmap_slide(data)
        self.prs.save(output_path)
        print(f"PPTX generated: {output_path}")

    def _add_title_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        color = self._get_color(data['grade'])

        # Brief Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        p = title.text_frame.paragraphs[0]
        p.text = "Security Posture: Executive Brief"
        p.font.size, p.font.bold, p.alignment = Pt(44), True, PP_ALIGN.CENTER

        # Grade Letter
        grade_box = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(5), Inches(2.5))
        p = grade_box.text_frame.paragraphs[0]
        p.text = data['grade']
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(180), True, color, PP_ALIGN.CENTER

        # Metadata
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        p = footer.text_frame.paragraphs[0]
        total_hours = data.get('total_effort_hours', 0)
        p.text = f"Report ID: {data.get('report_id', 'INTERNAL')} | Findings: {data['summary']['total_findings']} | Est. Effort: {total_hours}h"
        p.font.size, p.alignment = Pt(14), PP_ALIGN.CENTER

    def _add_matrix_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        plan = data['execution_plan']
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title.text_frame.text = "Security Explain Plan: Resource Allocation"
        
        # 4-tier system with NESTED_LOOP (Medium) restored
        rows = [
            ("FULL TABLE SCAN", plan['full_table_scans'], RGBColor(220, 53, 69), "IMMEDIATE", "Critical / CISA KEV"),
            ("INDEX RANGE SCAN", plan['index_scans'], RGBColor(253, 126, 20), "NEXT SPRINT", "High Severity"),
            ("NESTED LOOP", plan['nested_loops'], RGBColor(255, 193, 7), "SPRINT +1", "Medium Severity"),
            ("SEQUENTIAL READ", plan['low_priority'], RGBColor(108, 117, 125), "BACKLOG", "Low Severity")
        ]

        # Adjust spacing for 4 rows instead of 3
        row_height = 1.3
        start_y = 1.5
        
        for i, (label, section, color, priority, subtitle) in enumerate(rows):
            y = start_y + (i * row_height)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
            shape.line.color.rgb, shape.line.width = color, Pt(2)
            
            tf = shape.text_frame
            tf.word_wrap = True
            
            # Title line
            p1 = tf.paragraphs[0]
            p1.text = f"{label} ({priority})"
            p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(16), True, color
            
            # Subtitle line
            p2 = tf.add_paragraph()
            p2.text = f"{subtitle} | {section['count']} Findings | {section['estimated_hours']}h Effort"
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(100, 100, 100)

    def _add_critical_detail_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        items = data['execution_plan']['full_table_scans']['items']
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title.text_frame.text = "Critical Risk Path (Full Table Scans)"
        
        if not items:
            msg = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            msg.text_frame.text = "No Critical Findings Identified"
            return

        y = 1.6
        for item in items[:3]:  # Top 3 critical
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.4))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 245, 245)
            
            tf = box.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            kev = "[CISA KEV] " if item.get('cisa_kev') else ""
            raw_title = item.get('title', '')
            display_title = (raw_title[:60] + '...') if len(raw_title) > 63 else raw_title
            p1.text = f"{kev}{item.get('id')}: {display_title}"
            p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(14), True, RGBColor(220, 53, 69)
            
            p2 = tf.add_paragraph()
            fix = item.get('fixed_version') or 'Contact Vendor'
            p2.text = f"Package: {item.get('pkg_name')} | Fix: {fix} | Effort: {item.get('fix_effort_hours', '?')}h"
            p2.font.size = Pt(11)
            
            y += 1.6

    def _add_high_detail_slide(self, data):
        """New slide for high severity findings."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        items = data['execution_plan']['index_scans']['items']
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title.text_frame.text = "High Priority Queue (Index Range Scans)"
        
        if not items:
            msg = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            msg.text_frame.text = "No High Severity Findings"
            return

        y = 1.6
        for item in items[:3]:  # Top 3 high
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.4))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 252, 245)
            
            tf = box.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            raw_title = item.get('title', '')
            display_title = (raw_title[:60] + '...') if len(raw_title) > 63 else raw_title
            p1.text = f"{item.get('id')}: {display_title}"
            p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(14), True, RGBColor(253, 126, 20)
            
            p2 = tf.add_paragraph()
            fix = item.get('fixed_version') or 'TBD'
            p2.text = f"Package: {item.get('pkg_name')} | Fix: {fix} | Effort: {item.get('fix_effort_hours', '?')}h"
            p2.font.size = Pt(11)
            
            y += 1.6

    def _add_roadmap_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1)).text_frame.text = "Remediation Roadmap"
        
        plan = data['execution_plan']
        crit = plan['full_table_scans']['count']
        high = plan['index_scans']['count']
        med = plan['nested_loops']['count']
        low = plan['low_priority']['count']

        phases = [
            ("Phase 1: Emergency Response", f"Fix {crit} Critical (Full Table Scans)", RGBColor(220, 53, 69)),
            ("Phase 2: Sprint Commitment", f"Schedule {high} High Priority (Index Scans)", RGBColor(253, 126, 20)),
            ("Phase 3: Planned Maintenance", f"Queue {med} Medium (Nested Loops)", RGBColor(255, 193, 7)),
            ("Phase 4: Continuous Hardening", f"Review {low} Low Priority items", RGBColor(108, 117, 125))
        ]

        y = 1.5
        for label, detail, color in phases:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(11), Inches(1.0))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
            shape.line.color.rgb = color
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(14), True, color
            
            p2 = tf.add_paragraph()
            p2.text = detail
            p2.font.size = Pt(11)
            
            y += 1.2
