from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime

class PPTXGenerator:
    def __init__(self, master_pptx: Path = None):
        self.prs = Presentation(str(master_pptx)) if master_pptx and master_pptx.exists() else Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _get_color(self, grade: str) -> RGBColor:
        colors = {
            'A': RGBColor(40, 167, 69), 'B': RGBColor(108, 117, 125),
            'C': RGBColor(255, 193, 7), 'D': RGBColor(253, 126, 20),
            'F': RGBColor(220, 53, 69)
        }
        return colors.get(grade, RGBColor(0, 0, 0))

    def _ensure_data_structure(self, data: dict) -> dict:
        """Sanitizes input data for slide stability."""
        data.setdefault('grade', 'F')
        data.setdefault('generated_at', datetime.now().strftime('%Y-%m-%d %H:%M'))
        data.setdefault('summary', {'total_findings': 0})
        
        ep = data.setdefault('execution_plan', {})
        for section in ['full_table_scans', 'index_scans', 'low_priority']:
            ep.setdefault(section, {'count': 0, 'estimated_hours': 0, 'items': []})
        return data

    def generate_pptx(self, data: dict, output_path: str):
        data = self._ensure_data_structure(data)
        self._add_title_slide(data)
        self._add_matrix_slide(data)
        self._add_critical_detail_slide(data)
        self._add_roadmap_slide(data)
        self.prs.save(output_path)
        print(f"PPTX generated: {output_path}")

    def _add_title_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        color = self._get_color(data['grade'])

        # Brief Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        p = title.text_frame.paragraphs[0]
        p.text = "Security Posture: Executive Brief"
        p.font.size, p.font.bold, p.alignment = Pt(44), True, PP_ALIGN.CENTER

        # Grade Letter
        grade_box = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(5), Inches(2.5))
        p = grade_box.text_frame.paragraphs[0]
        p.text = data['grade']
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(180), True, color, PP_ALIGN.CENTER

        # Metadata
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        p = footer.text_frame.paragraphs[0]
        p.text = f"Report ID: {data.get('report_id', 'INTERNAL')} | Findings: {data['summary']['total_findings']}"
        p.font.size, p.alignment = Pt(14), PP_ALIGN.CENTER

    def _add_matrix_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        plan = data['execution_plan']
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title.text_frame.text = "Security Explain Plan: Resource Allocation"
        
        rows = [
            ("FULL TABLE SCAN", plan['full_table_scans'], RGBColor(220, 53, 69), "IMMEDIATE"),
            ("INDEX RANGE SCAN", plan['index_scans'], RGBColor(253, 126, 20), "NEXT SPRINT"),
            ("DEFERRED OPERATIONS", plan['low_priority'], RGBColor(108, 117, 125), "BACKLOG")
        ]

        for i, (label, section, color, priority) in enumerate(rows):
            y = 1.8 + (i * 1.7)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
            shape.line.color.rgb, shape.line.width = color, Pt(2)
            
            tf = shape.text_frame
            tf.paragraphs[0].text = f"{label} ({priority})"
            tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(18), True, color
            
            p2 = tf.add_paragraph()
            p2.text = f"{section['count']} Findings | {section['estimated_hours']}h Estimated Effort"
            p2.font.size = Pt(24)

    def _add_critical_detail_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        items = data['execution_plan']['full_table_scans']['items']
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title.text_frame.text = "Critical Risk Path"
        
        if not items:
            msg = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            msg.text_frame.text = "No Critical Findings Identified"
            return

        y = 1.6
        for item in items[:3]:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.4))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 245, 245)
            
            tf = box.text_frame
            p1 = tf.paragraphs[0]
            kev = "[CISA KEV] " if item.get('cisa_kev') else ""
            raw_title = item.get('title', '')
            title = (raw_title[:77] + '...') if len(raw_title) > 80 else raw_title
            p1.text = f"{kev}{item.get('id')}: {title}"
            p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(16), True, RGBColor(220, 53, 69)
            
            p2 = tf.add_paragraph()
            p2.text = f"Package: {item.get('pkg_name')} | Fix: {item.get('fixed_version') or 'Contact Vendor'}"
            p2.font.size = Pt(20)
            y += 1.6

    def _add_roadmap_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1)).text_frame.text = "Remediation Roadmap"
        
        # Simple timeline logic
        y = 1.8
        crit = data['execution_plan']['full_table_scans']['count']
        high = data['execution_plan']['index_scans']['count']

        for phase, detail in [
            ("Phase 1", f"Address {crit} Critical Findings"),
            ("Phase 2", f"Remediate {high} High-Risk Issues"),
            ("Phase 3", "Ongoing Monitoring & Hardening")
        ]:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(11), Inches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.text_frame.text = f"{phase}: {detail}"
            y += 1.5
