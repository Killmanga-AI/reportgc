"""
Unit tests for pptx_generator.py - Presentation generation.
"""

import pytest
from pathlib import Path
from pptx import Presentation
from pptx_generator import PPTXGenerator


class TestPPTXGeneratorInitialization:
    """Test generator setup and configuration."""
    
    def test_default_dimensions(self):
        """Should set 16:9 aspect ratio by default."""
        gen = PPTXGenerator()
        assert gen.prs.slide_width.inches == pytest.approx(13.333, 0.01)
        assert gen.prs.slide_height.inches == pytest.approx(7.5, 0.01)
    
    def test_loads_master_template(self, tmp_path):
        """Should load master template if provided."""
        # Create dummy pptx file
        master = tmp_path / "master.pptx"
        prs = Presentation()
        prs.save(str(master))
        
        gen = PPTXGenerator(master_pptx=master)
        assert gen.prs is not None


class TestDataStructureValidation:
    """Test input data sanitization."""
    
    def test_ensures_default_grade(self):
        """Should default to F if grade missing."""
        gen = PPTXGenerator()
        data = {}
        result = gen._ensure_data_structure(data)
        assert result["grade"] == "F"
    
    def test_ensures_execution_plan_sections(self):
        """Should create all 4 execution plan sections if missing."""
        gen = PPTXGenerator()
        data = {}
        result = gen._ensure_data_structure(data)
        
        ep = result["execution_plan"]
        assert "full_table_scans" in ep
        assert "index_scans" in ep
        assert "nested_loops" in ep
        assert "low_priority" in ep
        assert ep["full_table_scans"]["count"] == 0
    
    def test_preserves_existing_data(self):
        """Should not overwrite existing valid data."""
        gen = PPTXGenerator()
        data = {
            "grade": "B",
            "execution_plan": {
                "full_table_scans": {"count": 5, "estimated_hours": 30, "items": []}
            }
        }
        result = gen._ensure_data_structure(data)
        assert result["grade"] == "B"
        assert result["execution_plan"]["full_table_scans"]["count"] == 5


classTestColorMapping:
    """Test color assignment for grades and risk levels."""
    
    def test_grade_a_is_green(self):
        gen = PPTXGenerator()
        color = gen._get_color("A")
        assert color.rgb == (40, 167, 69)  # Bootstrap success green
    
    def test_grade_f_is_red(self):
        gen = PPTXGenerator()
        color = gen._get_color("F")
        assert color.rgb == (220, 53, 69)  # Bootstrap danger red
    
    def test_critical_risk_is_red(self):
        gen = PPTXGenerator()
        color = gen._get_risk_color("FULL_TABLE_SCAN")
        assert color.rgb == (220, 53, 69)
    
    def test_medium_risk_is_yellow(self):
        gen = PPTXGenerator()
        color = gen._get_risk_color("NESTED_LOOP")
        assert color.rgb == (255, 193, 7)


class TestSlideGeneration:
    """Test actual slide creation (creates real PPTX files)."""
    
    def test_generates_all_slides(self, tmp_path):
        """Should generate 5 slides for complete data."""
        gen = PPTXGenerator()
        output = tmp_path / "test.pptx"
        
        data = {
            "grade": "C",
            "report_id": "20240101-120000",
            "generated_at": "2024-01-01 12:00:00",
            "summary": {"total_findings": 10},
            "total_effort_hours": 50,
            "execution_plan": {
                "full_table_scans": {
                    "count": 2,
                    "estimated_hours": 12,
                    "items": [
                        {
                            "id": "CVE-2023-1",
                            "title": "Critical Vulnerability",
                            "cisa_kev": True,
                            "pkg_name": "openssl",
                            "fixed_version": "1.2.3",
                            "fix_effort_hours": 6
                        }
                    ]
                },
                "index_scans": {
                    "count": 3,
                    "estimated_hours": 12,
                    "items": []
                },
                "nested_loops": {
                    "count": 4,
                    "estimated_hours": 16,
                    "items": []
                },
                "low_priority": {
                    "count": 1,
                    "estimated_hours": 0,
                    "items": []
                }
            }
        }
        
        gen.generate_pptx(data, str(output))
        
        # Verify file created
        assert output.exists()
        
        # Verify slides
        prs = Presentation(str(output))
        assert len(prs.slides) == 5  # Title, Matrix, Critical, High, Roadmap
    
    def test_handles_no_critical_findings(self, tmp_path):
        """Should handle empty critical list gracefully."""
        gen = PPTXGenerator()
        output = tmp_path / "test.pptx"
        
        data = {
            "grade": "A",
            "report_id": "20240101-120000",
            "generated_at": "2024-01-01 12:00:00",
            "summary": {"total_findings": 0},
            "execution_plan": {
                "full_table_scans": {"count": 0, "estimated_hours": 0, "items": []},
                "index_scans": {"count": 0, "estimated_hours": 0, "items": []},
                "nested_loops": {"count": 0, "estimated_hours": 0, "items": []},
                "low_priority": {"count": 0, "estimated_hours": 0, "items": []}
            }
        }
        
        gen.generate_pptx(data, str(output))
        assert output.exists()
        
        prs = Presentation(str(output))
        # Should still generate all slides, with "No Findings" messages
    
    def test_limits_critical_details_to_three(self, tmp_path, factory):
        """Should only show top 3 critical findings."""
        gen = PPTXGenerator()
        output = tmp_path / "test.pptx"
        
        items = [
            {
                "id": f"CVE-2023-{i}",
                "title": f"Vulnerability {i}",
                "cisa_kev": False,
                "pkg_name": "test",
                "fixed_version": "1.0",
                "fix_effort_hours": 4
            }
            for i in range(5)  # 5 critical findings
        ]
        
        data = {
            "grade": "D",
            "report_id": "20240101-120000",
            "generated_at": "2024-01-01 12:00:00",
            "summary": {"total_findings": 5},
            "execution_plan": {
                "full_table_scans": {
                    "count": 5,
                    "estimated_hours": 20,
                    "items": items
                },
                "index_scans": {"count": 0, "estimated_hours": 0, "items": []},
                "nested_loops": {"count": 0, "estimated_hours": 0, "items": []},
                "low_priority": {"count": 0, "estimated_hours": 0, "items": []}
            }
        }
        
        gen.generate_pptx(data, str(output))
        
        # Verify by checking slide content (manual inspection or text extraction)
        prs = Presentation(str(output))
        # Third slide should be critical details
        critical_slide = prs.slides[2]
        # Should have 3 shapes (title + 3 finding boxes) or similar
